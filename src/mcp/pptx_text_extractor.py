"""
Extract plain text from PowerPoint (.pptx) files for cloud drive MCP reads.

Used by drive_file_extractor (Google Drive + OneDrive). Not registered in
DocumentLoader / KB indexing directly.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None  # type: ignore[misc, assignment]
    MSO_SHAPE_TYPE = None  # type: ignore[misc, assignment]
    PPTX_AVAILABLE = False
    logger.warning(
        "python-pptx not installed. Drive .pptx reads will fail until "
        "pip install python-pptx is applied on the API server."
    )


def _shape_text_parts(shape: Any) -> List[str]:
    """Collect text from a single shape (text frame, table, or group)."""
    parts: List[str] = []
    if not PPTX_AVAILABLE or MSO_SHAPE_TYPE is None:
        return parts

    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                parts.extend(_shape_text_parts(child))
            return parts

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            for row in table.rows:
                row_cells = []
                for cell in row.cells:
                    cell_text = (cell.text or "").strip()
                    if cell_text:
                        row_cells.append(cell_text)
                if row_cells:
                    parts.append(", ".join(row_cells))
            return parts

        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            frame = shape.text_frame
            for paragraph in frame.paragraphs:
                line = "".join(run.text for run in paragraph.runs).strip()
                if not line:
                    line = (paragraph.text or "").strip()
                if line:
                    parts.append(line)
    except Exception as exc:
        logger.debug("Skipping shape text extraction: %s", exc)

    return parts


def _slide_text(slide: Any, slide_number: int, include_notes: bool = True) -> str:
    """Build text block for one slide."""
    body_parts: List[str] = []
    for shape in slide.shapes:
        body_parts.extend(_shape_text_parts(shape))

    if include_notes:
        try:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                for paragraph in notes_slide.notes_text_frame.paragraphs:
                    note_line = (paragraph.text or "").strip()
                    if note_line:
                        body_parts.append(f"[Speaker note] {note_line}")
        except Exception as exc:
            logger.debug("No speaker notes on slide %s: %s", slide_number, exc)

    if not body_parts:
        return ""

    return f"--- Slide {slide_number} ---\n" + "\n".join(body_parts)


def extract_pptx_text(file_path: str, include_notes: bool = True) -> Dict[str, Any]:
    """
    Extract slide text from a .pptx file on disk.

    Returns:
        dict with keys: success (bool), full_text (str), slide_count (int),
        extraction_method (str), error (optional str)
    """
    path = Path(file_path)
    logger.info("extract_pptx_text: starting path=%s", path)

    if not PPTX_AVAILABLE or Presentation is None:
        return {
            "success": False,
            "full_text": "",
            "slide_count": 0,
            "extraction_method": "pptx",
            "error": "python-pptx not installed",
        }

    if not path.exists():
        logger.warning("extract_pptx_text: file not found path=%s", path)
        return {
            "success": False,
            "full_text": "",
            "slide_count": 0,
            "extraction_method": "pptx",
            "error": f"file not found: {path}",
        }

    if path.suffix.lower() not in {".pptx", ".ppt"}:
        logger.warning("extract_pptx_text: unexpected suffix=%s", path.suffix)

    try:
        presentation = Presentation(str(path))
        slide_blocks: List[str] = []

        for idx, slide in enumerate(presentation.slides, start=1):
            block = _slide_text(slide, idx, include_notes=include_notes)
            if block:
                slide_blocks.append(block)
            else:
                logger.debug("extract_pptx_text: slide %s has no extractable text", idx)

        full_text = "\n\n".join(slide_blocks).strip()
        slide_count = len(presentation.slides)

        if not full_text:
            logger.warning(
                "extract_pptx_text: no text extracted slide_count=%s path=%s",
                slide_count,
                path.name,
            )
            return {
                "success": False,
                "full_text": "",
                "slide_count": slide_count,
                "extraction_method": "pptx",
                "error": "no extractable text in slides (image-only deck?)",
            }

        logger.info(
            "extract_pptx_text: success path=%s slides=%s chars=%s",
            path.name,
            slide_count,
            len(full_text),
        )
        return {
            "success": True,
            "full_text": full_text,
            "slide_count": slide_count,
            "extraction_method": "pptx",
            "error": None,
        }
    except Exception as exc:
        logger.error("extract_pptx_text failed path=%s: %s", path, exc, exc_info=True)
        return {
            "success": False,
            "full_text": "",
            "slide_count": 0,
            "extraction_method": "pptx",
            "error": str(exc),
        }
